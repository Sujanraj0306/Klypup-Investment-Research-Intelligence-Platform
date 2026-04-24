"""Generate the Case Study Word doc and PowerPoint deck.

Run from repo root with the backend venv active:
    backend/.venv/bin/python scripts/generate_deliverables.py

Outputs:
    docs/Klypup_Case_Study.docx
    docs/Klypup_Case_Study.pptx

The PPT is text-heavy, narrative, and technical-depth-appropriate. It does
NOT embed screenshots — the presenter uses their own screenshots and video.
"""

from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PPTInches
from pptx.util import Pt as PPTPt


ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
DEMO_URL = "https://youtu.be/7bVMFYfFz_A"

# ── Minimal light theme ──────────────────────────────────────────────────
BG = PPTRGBColor(0xFA, 0xFA, 0xF7)
TEXT = PPTRGBColor(0x11, 0x1B, 0x30)
SUBTLE = PPTRGBColor(0x4A, 0x55, 0x6A)
ACCENT = PPTRGBColor(0x2D, 0x5B, 0xD8)
ACCENT_SOFT = PPTRGBColor(0xE3, 0xEC, 0xFF)
RULE = PPTRGBColor(0xD9, 0xDE, 0xE6)


# =========================================================================
# DOCX helpers
# =========================================================================


def _set_cell_bg(cell, hex_color: str):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tc_pr.append(shd)


def _h1(doc, text):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(0x11, 0x1B, 0x30)
    p.paragraph_format.space_after = Pt(6)


def _h2(doc, text):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = True
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(0x2D, 0x5B, 0xD8)
    p.paragraph_format.space_before = Pt(14)
    p.paragraph_format.space_after = Pt(4)


def _para(doc, text, *, size=11, bold=False, italic=False,
          color=(0x22, 0x2A, 0x3C)):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.size = Pt(size)
    run.bold = bold
    run.italic = italic
    run.font.color.rgb = RGBColor(*color)
    p.paragraph_format.space_after = Pt(6)
    return p


def _bullet(doc, text):
    p = doc.add_paragraph(style="List Bullet")
    if p.runs:
        run = p.runs[0]
        run.text = text
    else:
        run = p.add_run(text)
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x22, 0x2A, 0x3C)


def _stats_table(doc, rows):
    table = doc.add_table(rows=len(rows) + 1, cols=3)
    table.style = "Light Grid Accent 1"
    for i, h in enumerate(["What", "Stat", "Source"]):
        cell = table.rows[0].cells[i]
        cell.text = h
        for run in cell.paragraphs[0].runs:
            run.bold = True
            run.font.size = Pt(10.5)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _set_cell_bg(cell, "2D5BD8")
    for idx, (label, value, source) in enumerate(rows, start=1):
        c = table.rows[idx].cells
        c[0].text = label
        c[1].text = value
        c[2].text = source
        for cell in c:
            for p in cell.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)


def _demo_callout(doc):
    table = doc.add_table(rows=1, cols=1)
    table.autofit = True
    cell = table.rows[0].cells[0]
    _set_cell_bg(cell, "E3ECFF")
    para = cell.paragraphs[0]
    para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = para.add_run("▶  Watch the 10-minute demo first  →  ")
    r.bold = True
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor(0x11, 0x1B, 0x30)
    link = para.add_run(DEMO_URL)
    link.font.size = Pt(13)
    link.font.color.rgb = RGBColor(0x2D, 0x5B, 0xD8)
    link.bold = True
    link.underline = True
    doc.add_paragraph()


def build_docx():
    doc = Document()
    for section in doc.sections:
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.85)
        section.right_margin = Inches(0.85)

    # Cover
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = t.add_run("Klypup")
    tr.bold = True; tr.font.size = Pt(36); tr.font.color.rgb = RGBColor(0x11, 0x1B, 0x30)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = subtitle.add_run("Investment Research Intelligence — Business & Technical Case Study")
    sr.font.size = Pt(14); sr.font.color.rgb = RGBColor(0x4A, 0x55, 0x6A)
    subtitle.paragraph_format.space_after = Pt(16)

    _demo_callout(doc)

    # 1. Problem
    _h1(doc, "1.  The problem")
    _para(
        doc,
        "A sell-side research analyst at a mid-sized investment bank is asked "
        "on Monday morning: \"How is NVIDIA's Q3 looking versus AMD, and what "
        "are the biggest risks?\" What actually happens next is a two-to-five "
        "day copy-paste marathon — pulling prices from the Bloomberg Terminal, "
        "reading forty-plus news articles by hand, opening 10-Qs on SEC EDGAR, "
        "and stitching everything into a Word document. By the time the "
        "report lands, the market has already moved.",
    )
    _h2(doc, "Why it happens — the numbers")
    _stats_table(doc, [
        ("Share of analyst time on gathering vs analysis",
         "~65% on data collection + validation, only 35% on insight",
         "FP&A Trends Survey 2024"),
        ("Single Bloomberg Terminal seat",
         "$31,980 per seat per year (2026 list price)",
         "Bloomberg pricing disclosure, 2026"),
        ("Global spend on financial market data",
         "$44.3 B in 2024 → $49.2 B in 2025 (+6.4% YoY)",
         "Burton-Taylor via Finextra"),
        ("Productivity ceiling with AI assistance",
         "Up to 80% time savings on data collection, 50% on analysis",
         "Moody's Research Assistant pilot"),
        ("Front-office IB productivity projection",
         "27–35% gains by 2026",
         "Deloitte"),
    ])

    # 2. Customer
    _h1(doc, "2.  The customer")
    _para(
        doc,
        "Three buyer segments feel this pain every day. Sell-side research "
        "analysts at mid-market banks refresh coverage on fifteen to thirty "
        "companies, each refresh consuming two to five working days. "
        "Buy-side analysts at long-only funds screen fifty-plus tickers "
        "every week and need sentiment, filings context, and competitor "
        "framing fast. Family-office principals and retail investors track "
        "around twenty holdings but cannot justify the $30,000-per-year "
        "Bloomberg seat, so they fall back to Yahoo Finance and leave "
        "signals on the table.",
    )

    # 3. Solution
    _h1(doc, "3.  Our solution — in one sentence")
    _para(
        doc,
        "Klypup turns a plain-English question into a cited, structured "
        "research report in 20 to 45 seconds, grounded in real market data, "
        "news, and SEC filings, delivered in a browser on any device.",
        size=12, bold=True,
    )

    # 4. Product
    _h1(doc, "4.  Product walkthrough")
    _para(
        doc,
        "The following sections describe each surface of the product. "
        "Screenshots accompany the live presentation; please also watch the "
        f"10-minute demo at {DEMO_URL} for the complete flow.",
        italic=True, color=(0x4A, 0x55, 0x6A),
    )

    _h2(doc, "4.1  Authentication & workspace")
    _para(
        doc,
        "Users sign in with Google or email and are automatically placed "
        "into a personal organisation that scopes every request. Firebase "
        "issues the ID token; every backend endpoint verifies it, resolves "
        "the organisation id, and injects it as a typed dependency. "
        "Multi-tenant isolation is enforced at three independent layers: "
        "Firestore security rules on `/orgs/{orgId}/**`, a backend "
        "middleware that refuses requests without a resolved org, and "
        "Supabase row-level security on the semantic-search table.",
    )

    _h2(doc, "4.2  Dashboard — the morning brief")
    _para(
        doc,
        "The dashboard replaces a Bloomberg terminal's opening screen. It "
        "renders an eleven-sector SPDR ETF heatmap driven by Financial "
        "Modeling Prep's stable-tier snapshot endpoint, the day's top-five "
        "gainers and losers, and a personal watchlist with seven-day price "
        "sparklines. Every tile is populated by the same caching layer "
        "(five-minute TTL on quotes, fifteen on the heatmap) that the rest "
        "of the app uses, so the morning landing is free after the first "
        "analyst in the org loads it.",
    )

    _h2(doc, "4.3  Research — the agentic flow")
    _para(
        doc,
        "The user types any natural-language query. The backend opens a "
        "Server-Sent-Events stream and hands the query to the research "
        "agent, which runs on Gemini 2.5 Flash using Google's function "
        "calling API following the same agentic pattern as the Agent "
        "Development Kit. Gemini is first asked, \"which of these six "
        "tools do you need to answer this?\" — market data, news and "
        "sentiment, SEC filings via retrieval-augmented generation, social "
        "signals from Reddit and Google Trends, live Google-grounded web "
        "search, and an optional Playwright scraper. Only the tools Gemini "
        "selects are executed. Their results are fed back in a second "
        "Gemini turn that emits the final structured JSON report.",
    )
    _para(
        doc,
        "The frontend receives six section events in order — market, news, "
        "filings, social, synthesis, risks — and animates each one in as "
        "it arrives. The user never sees a thirty-second black box; they "
        "watch the agent work. Every claim in the synthesis section ends "
        "with a bracketed citation pointing at the source that produced "
        "it, so the analyst can click through to Financial Modeling Prep, "
        "a NewsAPI article, or an SEC filing.",
    )

    _h2(doc, "4.4  Chat follow-up — grounded conversation")
    _para(
        doc,
        "Once a report completes, a chat panel appears to the right of "
        "the content. The user can ask questions like \"which is the "
        "better buy right now?\" or \"explain the biggest risk in detail\". "
        "The backend streams the response token-by-token using the same "
        "SSE protocol; the system prompt injects the key metrics, "
        "findings, and synthesis summary from the report so the answer is "
        "grounded in the numbers the user is looking at rather than "
        "Gemini's general knowledge.",
    )

    _h2(doc, "4.5  Compare — two to four companies side by side")
    _para(
        doc,
        "The compare view runs market, news, and social tools in "
        "parallel for each company using asyncio.gather, computes metric "
        "winners across eight normalised axes (revenue growth, gross "
        "margins, profit margins, valuation, forward valuation, revenue, "
        "market cap, momentum), and asks Gemini to produce a structured "
        "comparison plus an investor-archetype match (growth buy, value "
        "buy, momentum buy, income buy, sentiment buy). The UI renders "
        "a six-axis radar, a sortable metrics table with CSV export, a "
        "news-sentiment bar chart, and the AI synthesis.",
    )

    _h2(doc, "4.6  Reports archive — semantic search")
    _para(
        doc,
        "Every completed report is written to Firestore as the canonical "
        "source of truth. A background task embeds the report's query, "
        "summary, and company list with Gemini's `gemini-embedding-001` "
        "model (3072-dimensional vectors) and upserts into a Supabase "
        "Postgres table protected by the `pgvector` extension. A custom "
        "`match_reports` RPC performs cosine-similarity search scoped to "
        "the calling organisation. The analyst can later type \"GPU "
        "competitors\" and find a report titled \"NVIDIA vs AMD\" even "
        "though those exact words are not in the title.",
    )

    _h2(doc, "4.7  Watchlist")
    _para(
        doc,
        "The watchlist is a drag-to-reorder list of tickers backed by a "
        "sub-collection under the user's organisation. Each row shows a "
        "live price, a seven-day close sparkline, a change-percent chip, "
        "and a cached sentiment badge. Adding a ticker validates it via "
        "the same Financial Modeling Prep fallback chain used by the "
        "research tool so Yahoo rate-limiting cannot block a legitimate "
        "symbol.",
    )

    # 5. How it works (technical)
    _h1(doc, "5.  How it works — technical architecture")
    _para(
        doc,
        "The frontend is a Vite + React 18 + TypeScript single-page app "
        "styled with Tailwind (dark and light themes driven by CSS "
        "variables) and animated with Framer Motion. Charts use Recharts "
        "and the sector heatmap uses D3. React Query manages server "
        "state and a Zustand store holds ephemeral UI state. The SSE "
        "parser is hand-rolled because the browser's built-in "
        "`EventSource` only supports GET; our streams are POST + ID "
        "token.",
    )
    _para(
        doc,
        "The backend is FastAPI 0.115 on Python 3.12, served by Uvicorn. "
        "Every route requires a Firebase ID token; the `CurrentOrgUser` "
        "dependency verifies the token, reads `/users/{uid}.defaultOrg` "
        "from Firestore, and injects the org id into route handlers. "
        "The research agent lives in `app/agent/research_agent.py` and "
        "owns three Gemini turns: dispatch, tool execution, and "
        "synthesis. Tools live in `app/agent/tools.py` and each "
        "degrades gracefully when its upstream credentials are missing.",
    )
    _para(
        doc,
        "Retrieval-augmented generation over SEC filings runs in-process "
        "using ChromaDB with cosine-similarity search. Filings are "
        "fetched from SEC EDGAR, chunked with tiktoken's cl100k encoder "
        "at 600 tokens with 100 overlap, embedded with "
        "`gemini-embedding-001`, and persisted. Semantic search over "
        "saved reports uses the same embedding model but stores vectors "
        "in Supabase Postgres with pgvector for scalability.",
    )
    _para(
        doc,
        "Market data has a three-provider cascade. Financial Modeling "
        "Prep's `/stable/quote` is primary (no IP throttling on the free "
        "tier); yfinance is a secondary for richer fundamentals when "
        "Yahoo isn't rate-limited; Alpha Vantage's GLOBAL_QUOTE is the "
        "last-resort fallback for ETF symbols. The sector heatmap uses "
        "FMP's `/stable/sector-performance-snapshot` endpoint, and the "
        "movers widget uses `/stable/biggest-gainers` and "
        "`/stable/biggest-losers`. A module-level `_silence_stderr` "
        "context manager suppresses yfinance's rate-limit noise during "
        "fallback attempts.",
    )

    # 6. Impact
    _h1(doc, "6.  Business impact")
    _para(
        doc,
        "For the sell-side analyst, Klypup collapses the two-to-five day "
        "report cycle into thirty to sixty seconds for a first draft. "
        "A team that used to ship ten to fifteen reports per month can "
        "ship more than a hundred at the same headcount without losing "
        "the citation discipline. For the firm, the same research team "
        "expands coverage five to eight times and gains a compliance "
        "audit trail — every report is saved with its tool calls, "
        "durations, and source URLs. For the retail or family-office "
        "user priced out of Bloomberg today, the free tier costs $0 per "
        "month and runs on a mobile browser.",
    )

    # 7. Cost
    _h1(doc, "7.  Cost model")
    _stats_table(doc, [
        ("Gemini 2.5 Flash (dispatcher + synthesis)", "~$0.60 per 1,000 queries", "Google API rate"),
        ("FMP + NewsAPI + GNews",                    "$0",                        "Free tiers"),
        ("Firestore + Firebase Auth",                "~$0.05 per 1,000 queries",  "Firebase free plan"),
        ("Supabase pgvector",                        "$0",                        "Free tier"),
        ("Render backend hosting",                   "$0",                        "750 hrs/mo free"),
        ("Total marginal cost",                      "≈ $0.65 per 1,000 queries", "Sub-dollar per thousand"),
    ])
    _para(
        doc,
        "For comparison, one Bloomberg Terminal seat is $31,980 per year. "
        "A Klypup power user running 100 queries a day for a full year "
        "incurs roughly $23.70 in marginal cost — the Terminal is "
        "approximately 1,000 times more expensive.",
        bold=True,
    )

    # 8. Robustness
    _h1(doc, "8.  Tech robustness")
    _para(
        doc,
        "No single upstream failure kills a research run. Each data "
        "domain has a fallback chain (prices: FMP → yfinance → Alpha "
        "Vantage; news: NewsAPI → GNews; web context: Gemini Google "
        "grounding). Multi-tenant isolation is defended at three "
        "independent layers (Firestore rules, backend middleware, "
        "Supabase RLS). Users see progress instead of a black-box wait "
        "because every step is a Server-Sent Event. Missing API keys "
        "disable only the affected tool; Gemini errors fall back to a "
        "deterministic heuristic orchestrator that still emits a "
        "structured report.",
    )

    # 9. Security
    _h1(doc, "9.  Security posture")
    _bullet(doc, "Firebase ID tokens verified on every request.")
    _bullet(doc, "No secrets in source; `.env` files gitignored; cloud secret store in production.")
    _bullet(doc, "Firestore and Supabase use at-rest encryption by default; TLS end-to-end.")
    _bullet(doc, "Supabase anon key is row-level-security locked; service-role key is backend-only.")
    _bullet(doc, "Every report records `tools_used` and `duration_ms` for end-to-end traceability.")

    # 10. Roadmap
    _h1(doc, "10.  Roadmap — next 90 days")
    _bullet(doc, "Live earnings-call transcripts ingested into the RAG corpus.")
    _bullet(doc, "Agent memory across queries (\"compare this to Tuesday's run\").")
    _bullet(doc, "Scheduled research runs — weekly auto-refresh of coverage reports.")
    _bullet(doc, "Slack and email delivery so reports reach analysts where they work.")
    _bullet(doc, "Model routing: Haiku for simple queries, Pro for deep synthesis — cost drops further.")
    _bullet(doc, "Admin analytics: firm-wide token spend, coverage heatmap.")

    # 11. Summary
    _h1(doc, "11.  Summary")
    _para(
        doc,
        "The analyst problem is time, not data. Data is a $44 billion "
        "industry; time is what Klypup gives back. The product is a "
        "browser-based research agent that answers natural-language "
        "queries in under a minute with sourced, structured reports. It "
        "runs on free or near-free infrastructure and costs roughly "
        "1,000 times less than a Bloomberg seat per user per year.",
    )
    _demo_callout(doc)

    out = DOCS / "Klypup_Case_Study.docx"
    doc.save(str(out))
    print(f"✓ wrote {out.relative_to(ROOT)}")


# =========================================================================
# PPTX — narrative, full-sentence, technical depth, minimal light theme
# =========================================================================


def _slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, text, *, left, top, width, height, size=18, bold=False,
              color=TEXT, align=PP_ALIGN.LEFT, italic=False, font="Inter",
              line_spacing=1.25):
    tb = slide.shapes.add_textbox(
        PPTInches(left), PPTInches(top), PPTInches(width), PPTInches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = PPTPt(size); run.font.bold = bold; run.font.italic = italic
    run.font.name = font; run.font.color.rgb = color
    return tb


def _add_paragraphs(slide, paragraphs, *, left, top, width, height, size=14,
                    color=TEXT, line_spacing=1.32, space_after=8):
    """Multiple paragraphs in one text frame for narrative slides."""
    tb = slide.shapes.add_textbox(
        PPTInches(left), PPTInches(top), PPTInches(width), PPTInches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, ptext in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = PPTPt(space_after)
        run = p.add_run()
        run.text = ptext
        run.font.size = PPTPt(size); run.font.name = "Inter"
        run.font.color.rgb = color
    return tb


def _add_rect(slide, *, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        PPTInches(left), PPTInches(top), PPTInches(width), PPTInches(height),
    )
    shape.adjustments[0] = 0.06
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line; shape.line.width = PPTPt(0.75)
    shape.shadow.inherit = False
    return shape


def _add_line(slide, *, left, top, length, color=RULE, weight=1.2):
    shape = slide.shapes.add_connector(
        1, PPTInches(left), PPTInches(top),
        PPTInches(left + length), PPTInches(top)
    )
    shape.line.color.rgb = color; shape.line.width = PPTPt(weight)
    return shape


def _add_footer(slide, n, total):
    _add_text(slide, "Klypup — Investment Research Intelligence",
              left=0.4, top=7.05, width=5, height=0.25, size=9, color=SUBTLE)
    _add_text(slide, DEMO_URL,
              left=5.2, top=7.05, width=4, height=0.25, size=9,
              color=ACCENT, align=PP_ALIGN.RIGHT)
    _add_text(slide, f"{n} / {total}",
              left=9.0, top=7.05, width=0.8, height=0.25, size=9,
              color=SUBTLE, align=PP_ALIGN.RIGHT)


def _title_slide(slide):
    _slide_background(slide, BG)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PPTInches(0), PPTInches(0),
        PPTInches(0.4), PPTInches(7.5)
    )
    accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background(); accent.shadow.inherit = False

    _add_text(slide, "Klypup", left=1.0, top=1.5, width=11, height=1.4,
              size=72, bold=True, color=TEXT)
    _add_text(slide, "Investment Research Intelligence",
              left=1.0, top=2.8, width=11, height=0.7, size=28, color=SUBTLE)
    _add_line(slide, left=1.0, top=3.75, length=2.5, color=ACCENT, weight=3)
    _add_paragraphs(
        slide,
        [
            "A browser-based research agent that turns a plain-English "
            "question into a sourced, structured, streaming report in "
            "20-45 seconds.",
            "Powered by Gemini 2.5 Flash function calling, six "
            "orchestrated data tools, Firestore for canonical state, and "
            "Supabase pgvector for semantic memory across past reports.",
        ],
        left=1.0, top=3.95, width=11.3, height=2.0, size=18, color=TEXT,
    )

    _add_rect(slide, left=1.0, top=6.1, width=7.5, height=0.75, fill=ACCENT_SOFT)
    _add_text(slide, f"▶  Watch the demo first  →  {DEMO_URL}",
              left=1.0, top=6.25, width=7.5, height=0.5,
              size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def _section_header_slide(slide, number, title, subtitle=""):
    _slide_background(slide, BG)
    _add_text(slide, number, left=0.8, top=2.4, width=2, height=1,
              size=96, bold=True, color=ACCENT)
    _add_text(slide, title, left=2.8, top=2.7, width=10, height=1.2,
              size=44, bold=True, color=TEXT)
    if subtitle:
        _add_text(slide, subtitle, left=2.8, top=4.0, width=10, height=1.4,
                  size=18, color=SUBTLE, line_spacing=1.35)


def _content_slide_header(slide, title, kicker=None):
    _slide_background(slide, BG)
    if kicker:
        _add_text(slide, kicker.upper(), left=0.5, top=0.35, width=12.3,
                  height=0.3, size=10, bold=True, color=ACCENT)
    _add_text(slide, title, left=0.5, top=0.65, width=12.3, height=0.8,
              size=30, bold=True, color=TEXT)
    _add_line(slide, left=0.5, top=1.45, length=12.3)


def _stat_card(slide, *, left, top, width, height, number, label, source):
    _add_rect(slide, left=left, top=top, width=width, height=height,
              fill=PPTRGBColor(0xFF, 0xFF, 0xFF), line=RULE)
    _add_text(slide, number, left=left + 0.2, top=top + 0.2,
              width=width - 0.4, height=0.9, size=34, bold=True, color=ACCENT)
    _add_text(slide, label, left=left + 0.2, top=top + 1.05,
              width=width - 0.4, height=height - 1.45, size=12, color=TEXT,
              line_spacing=1.25)
    _add_text(slide, source, left=left + 0.2, top=top + height - 0.45,
              width=width - 0.4, height=0.3, size=8.5, italic=True, color=SUBTLE)


def build_pptx():
    prs = Presentation()
    prs.slide_width = PPTInches(13.333)
    prs.slide_height = PPTInches(7.5)
    blank = prs.slide_layouts[6]

    steps = [
        "title", "agenda", "problem_header", "problem_stats",
        "customer", "solution", "tech_overview", "agent_flow",
        "research_feature", "chat_feature", "compare_feature",
        "reports_feature", "dashboard_watchlist", "impact", "cost",
        "robustness", "security", "roadmap", "closing",
    ]
    total = len(steps)
    slides = [prs.slides.add_slide(blank) for _ in steps]

    # 1. Title
    _title_slide(slides[0])

    # 2. Agenda
    s = slides[1]
    _content_slide_header(s, "Agenda", kicker="What we'll cover")
    _add_paragraphs(
        s,
        [
            "The business problem — analysts pay $32,000 per seat every year to copy-paste data for days, and the market doesn't wait for them.",
            "The customer segments — sell-side, buy-side, and the long tail of family offices and retail investors priced out of Bloomberg today.",
            "The Klypup solution — one sentence, then a guided walk-through of each major feature with the technology behind it.",
            "Technical architecture — how the Gemini 2.5 Flash research agent picks tools at runtime, streams sections to the browser over Server-Sent Events, and stores semantic memory across reports.",
            "Business impact, cost economics, robustness, and the 90-day roadmap.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=15, color=TEXT,
        space_after=10,
    )
    _add_footer(s, 2, total)

    # 3. Problem header
    _section_header_slide(
        slides[2], "01", "The problem",
        "Research analysts already own the tools and the data. What they "
        "don't own is the time. A two-to-five day report cycle means the "
        "market has moved before the analyst can deliver the answer."
    )
    _add_footer(slides[2], 3, total)

    # 4. Problem stats
    s = slides[3]
    _content_slide_header(s, "The numbers behind the pain", kicker="The problem")
    stats = [
        ("~65%",  "of analyst time is spent on data gathering and validation, not insight generation.", "FP&A Trends Survey 2024"),
        ("$31,980", "is the annual list price of a single Bloomberg Terminal seat in 2026.", "Bloomberg pricing, 2026"),
        ("$44.3 B", "was spent globally on financial market data in 2024; $49.2 B in 2025.", "Burton-Taylor / Finextra"),
        ("27–35%", "front-office productivity gains are projected for top IBs by 2026 with AI.", "Deloitte"),
        ("80%",   "time savings on data collection reported in the Moody's Research Assistant pilot.", "Microsoft 2025"),
        ("½",     "cycle time — Schroders analysts go from two weeks to one with AI tools.", "Schroders"),
    ]
    col_w = 3.9; col_gap = 0.15
    for i, (num, label, src) in enumerate(stats):
        col = i % 3; row = i // 3
        left = 0.5 + col * (col_w + col_gap)
        top = 1.8 + row * 2.6
        _stat_card(s, left=left, top=top, width=col_w, height=2.4,
                   number=num, label=label, source=src)
    _add_footer(s, 4, total)

    # 5. Customer
    s = slides[4]
    _content_slide_header(s, "Who feels this pain every day", kicker="The customer")
    _add_paragraphs(
        s,
        [
            "Sell-side research analysts at mid-market investment banks refresh coverage on 15-30 companies, with each refresh consuming two to five working days. They already pay $30,000 a year for their Bloomberg seat, and the bulk of their workday is still copy-paste.",
            "Buy-side analysts at long-only funds screen 50+ tickers a week and need sentiment, filings context, and competitor framing faster than they can read. The same seat cost applies, with the same data-gathering burden.",
            "Family office principals and retail investors track around 20 holdings each. They cannot justify a terminal licence, so they fall back to Yahoo Finance and piecemeal news — and systematically leave signals on the table because there is no unified research surface in their price range.",
            "All three segments want the same thing: a trustworthy, cited answer in under a minute, not at the end of a three-day build.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=14, color=TEXT,
        space_after=10,
    )
    _add_footer(s, 5, total)

    # 6. Solution
    s = slides[5]
    _content_slide_header(s, "What Klypup is", kicker="The solution")
    _add_text(
        s,
        "Klypup turns a plain-English question into a cited, structured "
        "research report in 20 to 45 seconds, grounded in live market "
        "data, news, and SEC filings, delivered in a browser on any "
        "device.",
        left=0.8, top=1.85, width=11.7, height=1.6, size=22, color=TEXT,
        line_spacing=1.4,
    )
    _add_paragraphs(
        s,
        [
            "The product is a single-page React application talking to a FastAPI backend. The backend owns a research agent powered by Gemini 2.5 Flash that picks its own tools at runtime, runs them, and streams a structured JSON report back to the browser in real time.",
            "Every claim in the report is cited to its source — Financial Modeling Prep, NewsAPI, an SEC filing, a Reddit thread, or a live Google search result. The analyst can click through and verify.",
        ],
        left=0.8, top=3.9, width=11.7, height=2.5, size=14, color=SUBTLE,
        space_after=10,
    )
    _add_rect(s, left=0.8, top=6.3, width=11.7, height=0.6, fill=ACCENT_SOFT)
    _add_text(s, f"▶  Full walk-through on video  →  {DEMO_URL}",
              left=0.8, top=6.42, width=11.7, height=0.4, size=14, bold=True,
              color=ACCENT, align=PP_ALIGN.CENTER)
    _add_footer(s, 6, total)

    # 7. Tech stack overview
    s = slides[6]
    _content_slide_header(s, "The technology behind it", kicker="Architecture")
    _add_paragraphs(
        s,
        [
            "The frontend is a Vite-built React 18 single-page application in TypeScript, styled with Tailwind, animated with Framer Motion, and charted with Recharts and D3. It talks to the backend over REST and Server-Sent Events.",
            "The backend is FastAPI on Python 3.12 with Uvicorn. It enforces Firebase Auth on every route, resolves the caller's organisation id, and dispatches to an agent that runs on Gemini 2.5 Flash using Google's function-calling API — the same agentic pattern as the Agent Development Kit (ADK) but implemented directly against Gemini for tighter control over streaming.",
            "Data lives in two places. Firestore holds users, organisations, watchlists, and the canonical copy of every saved report. Supabase Postgres with the pgvector extension holds 3072-dimensional embeddings of every report, enabling semantic search over an analyst's own history.",
            "Retrieval-augmented generation over SEC filings runs in-process with ChromaDB, using tiktoken's cl100k encoder to chunk filings at 600 tokens with 100 overlap before embedding with gemini-embedding-001.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=13, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 7, total)

    # 8. Agent flow — how it actually works
    s = slides[7]
    _content_slide_header(s, "How the agent actually works", kicker="Under the hood")
    _add_paragraphs(
        s,
        [
            "Turn one — Gemini dispatch. The backend hands Gemini the user's query along with six function declarations: get_market_data, get_news_and_sentiment, query_sec_filings, get_social_trends, web_search (Google-grounded), and web_scrape. Gemini is asked which of these it wants to call to answer the question. This is not a hardcoded pipeline — a question about news alone does not trigger the filings tool.",
            "Turn two — local tool execution. The backend runs the selected tools. Market data goes through a three-provider cascade (Financial Modeling Prep → yfinance → Alpha Vantage). News aggregates NewsAPI and GNews with TextBlob sentiment classification. SEC filings run a cosine-similarity search in ChromaDB. Social pulls Reddit (PRAW) and Google Trends (pytrends). Web search uses Gemini's native google_search grounding. Each tool logs its duration so the presenter can show the timeline in the live demo.",
            "Turn three — synthesis. The raw tool outputs are fed back to Gemini with the instruction to emit a structured JSON report with six named sections (market, news, filings, social, synthesis, risks) plus a chat_context block for follow-up questions. Gemini returns the JSON; the backend walks the object and emits a Server-Sent-Event per section. The browser animates each section in as it lands.",
            "After the stream completes, the report is saved to Firestore. A background task embeds its query + summary + companies with gemini-embedding-001 and upserts into Supabase pgvector for future semantic search.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=12.5, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 8, total)

    # 9. Research feature
    s = slides[8]
    _content_slide_header(s, "Research — the core surface", kicker="Feature  ·  1")
    _add_paragraphs(
        s,
        [
            "The research page accepts any natural-language query through a textarea with voice-input support via the Web Speech API and a ⌘/Ctrl+K keyboard shortcut. A lightweight frontend validator rejects queries under ten characters or missing a ticker or financial keyword, preventing nonsense hits on the agent.",
            "When the user submits, the frontend opens an SSE stream to POST /api/research/stream with their Firebase ID token. The agent-status panel on the left renders every event the backend emits: the dispatch step, each tool call with its duration, the synthesis step, and a complete event with the saved report id.",
            "On the right, the ReportStreamView component listens for section events and swaps in the matching structured component — MarketSection (company cards plus a 30-day Recharts area chart), NewsSection (article list with sentiment badges and a sentiment gauge SVG), FilingsSection (RAG passages with relevance scores), SocialSection (Google Trends line plus Reddit posts), SynthesisSection (Markdown with read-aloud), RisksSection (severity-colored grid).",
            "A chat panel appears automatically when the report finishes, grounded in the report's key metrics and findings. The user can ask \"which is the better buy right now?\" and the answer streams token-by-token while citing the same numbers they just saw.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=12.5, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 9, total)

    # 10. Chat follow-up
    s = slides[9]
    _content_slide_header(s, "Chat follow-up — grounded conversation", kicker="Feature  ·  2")
    _add_paragraphs(
        s,
        [
            "The chat panel uses a separate SSE endpoint, POST /api/chat/followup. The request body carries the user's question, the completed report as report_context, and the last ten messages of conversation history so Gemini can keep context across turns.",
            "The backend builds a system prompt that injects the report's company list, price and valuation metrics, synthesis summary, and key findings. The prompt tells Gemini it is a financial research assistant — not a generic chatbot — and instructs it to answer with specific numbers, not refuse with compliance disclaimers.",
            "Gemini's response streams back chunk-by-chunk via Server-Sent Events so the user sees a typewriter effect. Each assistant message is rendered with react-markdown (bold key terms, bulleted lists) and every message has a read-aloud button that pipes the text into the Web Speech API's SpeechSynthesisUtterance.",
            "The entire chat is ephemeral in the page — it lives only in React state. If the analyst reloads, they return to the report with a fresh greeting. This keeps the compliance footprint small and avoids paying Gemini tokens for conversations the user abandoned.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=12.5, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 10, total)

    # 11. Compare
    s = slides[10]
    _content_slide_header(s, "Compare — two to four companies side by side", kicker="Feature  ·  3")
    _add_paragraphs(
        s,
        [
            "The compare view runs market, news, and social tools in parallel for each selected company using asyncio.gather in a single request. This returns in roughly the same time as a single-company research call — the win is strictly analyst throughput, not compute cost.",
            "The backend then computes per-metric winners across eight normalised axes: revenue growth, gross margins, profit margins, P/E valuation, forward P/E, total revenue, market cap, and momentum (one-day change percent). News sentiment is bucketed into a winner by the highest 0-100 sentiment score.",
            "A final Gemini turn is asked to produce a structured comparison object: a 3-4 paragraph narrative, an investor-archetype match (growth buy, value buy, momentum buy, income buy, sentiment buy), and 3-5 key differentiators. If Gemini fails for any reason, a deterministic heuristic generates an equivalent object so the comparison always renders.",
            "The UI renders a 6-axis radar chart (Recharts PolarGrid), a sortable metrics table with CSV export, a news-sentiment bar chart, and the AI synthesis — all fed by the same SSE stream that streams sections one at a time.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=12.5, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 11, total)

    # 12. Reports archive
    s = slides[11]
    _content_slide_header(s, "Reports archive — semantic memory", kicker="Feature  ·  4")
    _add_paragraphs(
        s,
        [
            "Every completed report is written to Firestore under /orgs/{orgId}/reports/{reportId} as the canonical source of truth. In parallel — fire-and-forget so the user never waits — a background task embeds the report's query, summary, and company list with Gemini's gemini-embedding-001 model (3072-dimensional output vectors) and upserts into a Supabase Postgres table.",
            "The table is protected by the pgvector extension and has a custom match_reports RPC that performs cosine-similarity search scoped to the caller's organisation. Row-level security ensures an org can only ever see its own rows, even if the anon key is leaked.",
            "In the UI, the analyst searches by meaning. Typing \"GPU competitors\" returns an NVIDIA-vs-AMD report even if those exact words do not appear in its title. Typing \"2026 outlook\" surfaces reports about forecasts across multiple companies. The analyst can tag reports (\"Q3 earnings\", \"competitor analysis\"), filter by tag, sort by recency or duration, and open any report back into the full structured view plus its chat panel.",
            "Reports are permanent — they live in Firestore under the organisation's path — but can be deleted. Deletion cleans both Firestore and the corresponding Supabase row.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=12.5, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 12, total)

    # 13. Dashboard + Watchlist
    s = slides[12]
    _content_slide_header(s, "Dashboard and watchlist", kicker="Feature  ·  5 + 6")
    _add_paragraphs(
        s,
        [
            "The dashboard is the analyst's morning brief. It renders an eleven-sector SPDR ETF heatmap (XLK, XLF, XLV, XLE, XLI, XLY, XLP, XLU, XLB, XLRE, XLC) driven by Financial Modeling Prep's /stable/sector-performance-snapshot endpoint. Below the heatmap sit the day's top-five gainers and losers pulled from /stable/biggest-gainers and /stable/biggest-losers, and the user's personal watchlist with seven-day sparklines.",
            "Every tile flows through the same backend cache layer — a tiny TTL dict with five-minute quotes, fifteen-minute sector heatmap, ten-minute movers, and negative-caching on failure so rate-limited upstreams don't thrash. The first analyst in the organisation to land on the dashboard warms the cache for everyone else.",
            "The watchlist is a drag-to-reorder list of tickers backed by a Firestore sub-collection under the user's organisation. Adding a ticker validates it through the FMP → yfinance → known-symbols cascade used elsewhere, so a Yahoo rate-limit cannot block a legitimate symbol like AAPL. Each row shows a live price, seven-day close sparkline, change-percent chip, and a cached sentiment badge.",
            "Both surfaces render in roughly 300ms after cache warm-up, with React Query handling background refetching so the UI always reflects sub-5-minute-old data.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=12.5, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 13, total)

    # 14. Business impact
    s = slides[13]
    _content_slide_header(s, "Business impact", kicker="Outcomes")
    cards = [
        ("For the sell-side analyst",
         "A two-to-five day report cycle collapses into thirty to sixty seconds for a first draft. A team that used to ship ten to fifteen reports per month can ship over one hundred at the same headcount without losing citation discipline. Every claim in every section is grounded in a source URL."),
        ("For the firm",
         "The same research team expands coverage five to eight times. A unified app replaces a meaningful share of what used to require a Bloomberg seat — live prices, news, filings, sentiment — so seat attrition becomes a real cost lever. Every report keeps a tools_used audit trail for compliance."),
        ("For the priced-out user",
         "Zero dollars per month on the free tier. Runs in any modern browser, including mobile. No Bloomberg licence, no Refinitiv subscription, no analyst headcount — just type a question and read a sourced answer. Compound interest on decisions that were previously out of reach."),
    ]
    col_w = 4.0; gap = 0.25
    total_w = 3 * col_w + 2 * gap
    start = (13.333 - total_w) / 2
    for i, (title, body) in enumerate(cards):
        left = start + i * (col_w + gap)
        _add_rect(s, left=left, top=1.85, width=col_w, height=4.9,
                  fill=PPTRGBColor(0xFF, 0xFF, 0xFF), line=RULE)
        _add_text(s, title, left=left + 0.25, top=2.05, width=col_w - 0.5,
                  height=0.8, size=16, bold=True, color=ACCENT, line_spacing=1.15)
        _add_text(s, body, left=left + 0.25, top=2.95, width=col_w - 0.5,
                  height=3.8, size=12, color=TEXT, line_spacing=1.35)
    _add_footer(s, 14, total)

    # 15. Cost
    s = slides[14]
    _content_slide_header(s, "Cost model — honest and auditable", kicker="Economics")
    _add_paragraphs(
        s,
        [
            "The marginal cost of running one thousand research queries on Klypup sits at roughly sixty-five cents. Almost all of it is Gemini 2.5 Flash tokens — the dispatcher turn plus the synthesis turn, averaging 30,000 input tokens and 4,000 output tokens per query at published rates.",
            "Financial Modeling Prep, NewsAPI, and GNews all provide sufficient free-tier coverage for demo workloads. Firestore and Firebase Auth reads and writes are within the free plan. Supabase pgvector with our data volume is free. Render hosts the backend for free with 750 hours per month.",
            "A single Bloomberg Terminal seat is $31,980 per year. A Klypup power user running one hundred queries a day for a full year incurs roughly $23.70 in marginal infrastructure cost — the Terminal is approximately 1,000 times more expensive per user per year.",
            "Scaling into a paid tier is still cheap. A $99-per-month Team plan at 80% gross margin supports a 25-analyst buy-side fund — one-fifteenth the cost of replacing their existing seat licences.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=13, color=TEXT,
        space_after=10,
    )
    _add_footer(s, 15, total)

    # 16. Robustness
    s = slides[15]
    _content_slide_header(s, "Why it stays reliable", kicker="Tech robustness")
    _add_paragraphs(
        s,
        [
            "No single upstream provider can kill a research run. Prices cascade through FMP → yfinance → Alpha Vantage. News aggregates NewsAPI plus GNews. Web context uses Gemini's google_search grounding, and if Gemini itself fails, a deterministic heuristic orchestrator still emits a structured report from whatever tool data came back.",
            "Multi-tenant isolation is defended at three independent layers: Firestore path-based security rules scoped to request.auth.uid, a backend CurrentOrgUser middleware that refuses requests without a resolved org, and Supabase row-level security on the semantic-search table. Any one of the three can fail and the other two still prevent a data leak.",
            "The user always sees progress, not a black box. Server-Sent Events push each tool start, tool completion, and section into the browser as soon as they land. A safety net in the React hook recovers even if an event is dropped mid-stream — it refetches the latest saved report so the UI never hangs.",
            "Missing API keys disable only the affected tool; Gemini errors fall back to the heuristic path; Yahoo rate-limit noise is suppressed at the stderr level so the production logs stay readable.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=12.5, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 16, total)

    # 17. Security
    s = slides[16]
    _content_slide_header(s, "Security posture", kicker="Trust")
    _add_paragraphs(
        s,
        [
            "Authentication is handled by Firebase Auth — Google SSO plus email/password. Every backend route is wrapped in a dependency that verifies the ID token, looks up the user's default organisation, and injects the typed org id. No route ever reads an org id from the request body, which closes the simplest multi-tenant exfiltration vector.",
            "Secrets are never checked in. The .env files are gitignored at the repo root, a pre-commit scan searches the staged diff for API-key-shaped strings before every commit, and production secrets live in the cloud provider's secret store — Render environment variables, never the container image.",
            "Data is encrypted at rest on Firestore and Supabase (both platform defaults) and transits over TLS end-to-end. The Supabase anon key shipped in the frontend is row-level-security locked — an attacker who steals it still sees zero rows. The service-role key exists only in the backend's memory.",
            "Every report records tools_used and duration_ms, giving compliance officers an audit trail from any answer back to the specific data calls that produced it.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=12.5, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 17, total)

    # 18. Roadmap
    s = slides[17]
    _content_slide_header(s, "Roadmap — next 90 days", kicker="Where this goes")
    _add_paragraphs(
        s,
        [
            "Live earnings-call transcripts ingested into the RAG corpus alongside 10-Ks and 10-Qs, so questions about forward guidance can be answered against the most recent management commentary rather than the last annual report.",
            "Agent memory across queries. A user can say \"compare this to Tuesday's run\" and the agent will hydrate the prior report, diff the numbers, and synthesise a change narrative without forcing the user to manually open both reports.",
            "Scheduled research runs — a user marks a coverage list and every Monday morning at 8am the agent re-runs each report, diffs against last week, and delivers the output via Slack or email so the analyst wakes up to a pre-populated workspace.",
            "Model routing. Simple quote lookups and rephrasing go through Haiku-class models for one-tenth the cost; deep synthesis sticks with Gemini 2.5 Flash or upgrades to Pro for complex multi-company comparisons. A per-query cost log feeds admin analytics.",
            "Admin analytics for firms — a dashboard of token spend per analyst, coverage heatmap by sector, and query volume trends so research directors can see where the team's attention is actually going.",
        ],
        left=0.8, top=1.85, width=11.7, height=5.0, size=12.5, color=TEXT,
        space_after=8,
    )
    _add_footer(s, 18, total)

    # 19. Closing
    s = slides[18]
    _slide_background(s, BG)
    _add_text(s, "Thank you", left=0.8, top=1.8, width=11.8, height=1.2,
              size=60, bold=True, color=TEXT)
    _add_line(s, left=0.8, top=3.1, length=3, color=ACCENT, weight=3)
    _add_text(s, "Klypup — Investment Research Intelligence",
              left=0.8, top=3.3, width=11.8, height=0.6, size=22, color=SUBTLE)
    _add_text(
        s,
        "Natural-language research agent · Gemini 2.5 Flash function calling · "
        "six orchestrated tools · Firestore + Supabase pgvector memory · "
        "Server-Sent-Event streaming · free-tier deployable.",
        left=0.8, top=4.1, width=11.8, height=1.5, size=14, color=TEXT,
        line_spacing=1.4,
    )
    _add_rect(s, left=0.8, top=5.8, width=11.8, height=0.9, fill=ACCENT_SOFT)
    _add_text(s, "▶  Watch the demo",
              left=0.8, top=5.88, width=11.8, height=0.45, size=14, bold=True,
              color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(s, DEMO_URL,
              left=0.8, top=6.25, width=11.8, height=0.5, size=20, bold=True,
              color=TEXT, align=PP_ALIGN.CENTER)
    _add_text(
        s,
        "github.com/Sujanraj0306/Klypup-Investment-Research-Intelligence-Platform",
        left=0.8, top=7.0, width=11.8, height=0.3, size=11, color=SUBTLE,
        align=PP_ALIGN.CENTER,
    )

    out = DOCS / "Klypup_Case_Study.pptx"
    prs.save(str(out))
    print(f"✓ wrote {out.relative_to(ROOT)}  ({total} slides)")


if __name__ == "__main__":
    DOCS.mkdir(parents=True, exist_ok=True)
    build_docx()
    build_pptx()
    print("\nDeliverables in docs/:")
    for p in sorted(DOCS.glob("Klypup_Case_Study*")):
        print(f"  · {p.relative_to(ROOT)}  ({p.stat().st_size // 1024} KB)")
